"""
Create an animated PowerPoint presentation for the Migration Analysis Tool
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    """Add a title slide with Bosch colors"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Red strip at top (Bosch red)
    red_strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.2)
    )
    red_strip.fill.solid()
    red_strip.fill.fore_color.rgb = RGBColor(230, 0, 0)  # Bosch red
    red_strip.line.color.rgb = RGBColor(230, 0, 0)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)  # Bosch dark blue
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(102, 102, 102)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Red strip at top
    red_strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.2)
    )
    red_strip.fill.solid()
    red_strip.fill.fore_color.rgb = RGBColor(230, 0, 0)
    red_strip.line.color.rgb = RGBColor(230, 0, 0)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.space_before = Pt(12)
        p.level = 0
        
        # Check for sub-bullets (indented with -)
        if item.strip().startswith('-'):
            p.text = item.replace('-', '', 1).strip()
            p.level = 1
            p.font.size = Pt(18)
    
    return slide

def add_comparison_slide(prs):
    """Add time savings comparison slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Red strip
    red_strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.2)
    )
    red_strip.fill.solid()
    red_strip.fill.fore_color.rgb = RGBColor(230, 0, 0)
    red_strip.line.color.rgb = RGBColor(230, 0, 0)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "⚡ Time Savings Comparison"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Table header
    headers = ["Scenario", "Before Optimization", "After Optimization", "Speedup"]
    rows = [
        ["50 Files Processing", "2-3 minutes", "15-20 seconds", "10x faster"],
        ["200 Files Processing", "10-15 minutes", "1-2 minutes", "10x faster"],
        ["Single File Diff", "2-3 seconds", "0.5-1 second", "3x faster"],
    ]
    
    # Create table
    table_width = Inches(8.5)
    table_height = Inches(3)
    left = Inches(0.75)
    top = Inches(2.5)
    
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, table_width, table_height).table
    
    # Style header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Fill data rows
    for row_idx, row_data in enumerate(rows, start=1):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(14)
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(234, 243, 251)
    
    return slide

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, 
                   "Migration Analysis Tool",
                   "Automated Code Migration Analysis with AI Support")
    
    # Slide 2: Overview
    add_content_slide(prs, "🎯 Tool Overview", [
        "Comprehensive solution for analyzing code migration differences",
        "Supports multiple comparison modes and platforms",
        "RTC/ALM integration for enterprise workflows",
        "AI-powered intelligent merge suggestions",
        "Parallel processing for high-performance analysis",
        "Detailed HTML and Excel reporting"
    ])
    
    # Slide 3: Comparison Modes
    add_content_slide(prs, "📊 Comparison Modes", [
        "📁 Offline → Offline (Folder/ZIP Comparison)",
        "- Compare local directories or ZIP archives",
        "- Manual and automatic file mapping",
        "- Perfect for standalone analysis",
        "",
        "☁️ Online → Online (RTC Snapshot Comparison)",
        "- Compare live RTC snapshots directly",
        "- Automatic changeset tracking",
        "- Work item integration",
        "",
        "🔄 Online → Offline (Hybrid Mode)",
        "- Compare RTC snapshot with local folder",
        "- Best for validating local changes"
    ])
    
    # Slide 4: Key Features - Part 1
    add_content_slide(prs, "✨ Key Features - Analysis", [
        "🔍 Interface Analysis",
        "- Analyze header file interfaces (.h files)",
        "- Compare platform vs project interfaces",
        "- Identify API changes and dependencies",
        "",
        "🔗 Platform Dependency Analysis",
        "- Visualize platform-specific dependencies",
        "- Identify cross-module dependencies",
        "- Generate dependency reports",
        "",
        "📝 AI-Powered Smart Merge",
        "- Google Gemini AI integration",
        "- Intelligent merge conflict resolution",
        "- Context-aware suggestions"
    ])
    
    # Slide 5: Key Features - Part 2
    add_content_slide(prs, "✨ Key Features - Integration", [
        "🔐 RTC/ALM Integration",
        "- Direct RTC server connectivity",
        "- Fetch changesets and work items",
        "- Automatic snapshot fetching",
        "- Component-level analysis",
        "",
        "📊 Comprehensive Reporting",
        "- Detailed HTML diff reports",
        "- Excel reports with statistics",
        "- Color-coded change visualization",
        "- Export and share capabilities",
        "",
        "💬 AI Comparison Assistant",
        "- Ask questions about changes",
        "- Get intelligent explanations",
        "- Context-aware insights"
    ])
    
    # Slide 6: Time Savings
    add_comparison_slide(prs)
    
    # Slide 7: Optimization Details
    add_content_slide(prs, "🚀 Performance Optimizations", [
        "⚡ Parallel Processing (10x Speedup)",
        "- Process 10 files simultaneously",
        "- Concurrent file downloads and analysis",
        "- Optimized thread pool management",
        "",
        "💾 Memory-Based Diff Generation (3x Faster)",
        "- Eliminate temporary file I/O",
        "- In-memory diff computation",
        "- Direct HTML output generation",
        "",
        "📏 Smart File Filtering",
        "- Skip files larger than 500KB for HTML diffs",
        "- Prevent hanging on massive files",
        "- Focus on analyzable content"
    ])
    
    # Slide 8: Use Cases
    add_content_slide(prs, "💼 Use Cases", [
        "Platform Migration Projects",
        "- Compare old vs new platform files",
        "- Identify breaking changes",
        "- Track migration progress",
        "",
        "Release Validation",
        "- Verify release snapshots",
        "- Compare release candidates",
        "- Ensure code consistency",
        "",
        "Merge Conflict Resolution",
        "- AI-powered merge suggestions",
        "- Visualize conflicting changes",
        "- Streamline merge process",
        "",
        "Code Review & Audit",
        "- Detailed change analysis",
        "- Interface compatibility checks",
        "- Dependency impact assessment"
    ])
    
    # Slide 9: Benefits Summary
    add_content_slide(prs, "✅ Benefits Summary", [
        "⏱️ Massive Time Savings",
        "- 10x faster file processing",
        "- Automated changeset tracking",
        "- Quick interface analysis",
        "",
        "🎯 Accuracy & Completeness",
        "- Comprehensive change detection",
        "- Comment-only change identification",
        "- Detailed statistical reports",
        "",
        "🤝 Team Collaboration",
        "- Shareable HTML reports",
        "- Excel exports for management",
        "- Work item integration",
        "",
        "🔧 Easy to Use",
        "- Intuitive GUI interface",
        "- Multiple comparison modes",
        "- Real-time progress tracking"
    ])
    
    # Slide 10: Coverage Summary
    add_content_slide(prs, "📋 Tool Coverage", [
        "✓ File Comparison & Diff Generation",
        "✓ RTC Snapshot Management",
        "✓ Changeset & Work Item Tracking",
        "✓ Interface Analysis (Header Files)",
        "✓ Platform Dependency Visualization",
        "✓ AI-Powered Merge Suggestions",
        "✓ Excel & HTML Reporting",
        "✓ Parallel Processing Engine",
        "✓ Comment Detection",
        "✓ ZIP Archive Support",
        "✓ Credential Management",
        "✓ Real-Time Progress Updates"
    ])
    
    # Slide 11: Getting Started
    add_content_slide(prs, "🚀 Getting Started", [
        "1️⃣ Launch the Application",
        "- Run: python main.py",
        "- Or use the executable: Migration_Analysis.exe",
        "",
        "2️⃣ Choose Comparison Mode",
        "- Select: Offline, Online, or Hybrid mode",
        "- Configure inputs (folders, snapshots, etc.)",
        "",
        "3️⃣ Connect to RTC (Optional)",
        "- Enter credentials for RTC features",
        "- Enable changeset tracking",
        "",
        "4️⃣ Run Analysis & View Reports",
        "- Start comparison",
        "- View results in Excel/HTML",
        "- Ask AI assistant for insights"
    ])
    
    # Slide 12: Thank You
    add_title_slide(prs,
                   "Thank You!",
                   "Migration Analysis Tool - Bosch Engineering")
    
    return prs

if __name__ == "__main__":
    print("Creating Migration Analysis Tool presentation...")
    prs = create_presentation()
    output_path = "Migration_Analysis_Tool_Presentation.pptx"
    prs.save(output_path)
    print(f"✓ Presentation created: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")
    print("\nPresentation includes:")
    print("  • Tool overview and features")
    print("  • Comparison modes (Offline/Online/Hybrid)")
    print("  • Time savings analysis (10x faster)")
    print("  • Performance optimizations")
    print("  • Use cases and benefits")
    print("  • Complete tool coverage")
    print("  • Getting started guide")
