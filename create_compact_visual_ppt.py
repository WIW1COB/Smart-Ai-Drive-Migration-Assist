"""
Create a compact, highly attractive 5-slide presentation with flow diagrams
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_stunning_title_slide(prs):
    """Create a stunning title slide with visual elements"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Gradient background using overlapping shapes
    for i in range(5):
        alpha = 255 - (i * 40)
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(i * 1.5),
            Inches(10), Inches(2)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(234, 243, 251)
        bg.fill.transparency = i * 0.15
        bg.line.fill.background()
    
    # Large decorative circles
    circles = [
        (Inches(-1), Inches(-0.5), Inches(3), Inches(3), RGBColor(230, 0, 0), 0.15),
        (Inches(8), Inches(5), Inches(3), Inches(3), RGBColor(0, 51, 102), 0.15),
        (Inches(7), Inches(0.5), Inches(2), Inches(2), RGBColor(255, 140, 0), 0.2),
        (Inches(0.5), Inches(6), Inches(1.5), Inches(1.5), RGBColor(0, 153, 51), 0.2),
    ]
    
    for left, top, width, height, color, trans in circles:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.fill.transparency = trans
        circle.line.fill.background()
    
    # Red top strip
    red_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.4))
    red_strip.fill.solid()
    red_strip.fill.fore_color.rgb = RGBColor(230, 0, 0)
    red_strip.line.fill.background()
    
    # Central icon box - large and prominent
    icon_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(2),
        Inches(3), Inches(2)
    )
    icon_bg.fill.solid()
    icon_bg.fill.fore_color.rgb = RGBColor(0, 51, 102)
    icon_bg.line.width = Pt(5)
    icon_bg.line.color.rgb = RGBColor(230, 0, 0)
    icon_bg.shadow.inherit = False
    
    # Large icons
    icon_text = icon_bg.text_frame
    icon_text.text = "⚙️\n🚀"
    icon_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in icon_text.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(55)
    
    # Main title - bold and large
    title = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(1))
    tf = title.text_frame
    tf.text = "Migration Analysis Tool"
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle with key benefit
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(5.4), Inches(8), Inches(0.8))
    stf = subtitle.text_frame
    stf.text = "⚡ 10x Faster Code Migration Analysis with AI"
    p = stf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(230, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    
    # Feature badges at bottom
    badges = ["☁️ RTC Integration", "🤖 AI-Powered", "📊 Smart Reports"]
    x_positions = [1.5, 3.8, 6.1]
    
    for text, x in zip(badges, x_positions):
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(6.5),
            Inches(2.2), Inches(0.6)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(255, 215, 0)
        badge.line.width = Pt(2)
        badge.line.color.rgb = RGBColor(255, 140, 0)
        
        bt = badge.text_frame
        bt.text = text
        bt.paragraphs[0].font.size = Pt(14)
        bt.paragraphs[0].font.bold = True
        bt.paragraphs[0].alignment = PP_ALIGN.CENTER
        bt.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide

def add_overview_flow_slide(prs):
    """Slide 2: Overview with comprehensive flow diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(250, 252, 255)
    bg.line.fill.background()
    
    # Red strip
    red_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.3))
    red_strip.fill.solid()
    red_strip.fill.fore_color.rgb = RGBColor(230, 0, 0)
    red_strip.line.fill.background()
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(9), Inches(0.7))
    tf = title.text_frame
    tf.text = "🎯 Complete Migration Analysis Workflow"
    p = tf.paragraphs[0]
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # FLOW DIAGRAM - Complete Process
    
    # Step 1: Input Sources (Left side - 3 options)
    input_y_start = 1.6
    inputs = [
        ("📁 Folders/ZIPs", RGBColor(200, 230, 255)),
        ("☁️ RTC Snapshots", RGBColor(255, 230, 200)),
        ("🔄 Hybrid Mode", RGBColor(220, 255, 220))
    ]
    
    input_boxes = []
    for i, (text, color) in enumerate(inputs):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.4), Inches(input_y_start + i * 0.9),
            Inches(1.8), Inches(0.7)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.width = Pt(2)
        box.line.color.rgb = RGBColor(0, 51, 102)
        
        bt = box.text_frame
        bt.text = text
        bt.paragraphs[0].font.size = Pt(14)
        bt.paragraphs[0].font.bold = True
        bt.paragraphs[0].alignment = PP_ALIGN.CENTER
        bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        input_boxes.append((box, i))
    
    # Converging arrows from inputs to processing
    arrow_target_x = 2.4
    arrow_target_y = 2.8
    
    for i in range(3):
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(2.3), Inches(input_y_start + 0.2 + i * 0.9),
            Inches(0.9), Inches(0.3)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(230, 0, 0)
        arrow.line.fill.background()
    
    # Step 2: Processing Engine (Center)
    engine = slide.shapes.add_shape(
        MSO_SHAPE.FLOWCHART_PROCESS,
        Inches(3.2), Inches(2.3),
        Inches(1.8), Inches(1.4)
    )
    engine.fill.solid()
    engine.fill.fore_color.rgb = RGBColor(0, 51, 102)
    engine.line.width = Pt(3)
    engine.line.color.rgb = RGBColor(230, 0, 0)
    
    et = engine.text_frame
    et.text = "⚙️\nAnalysis\nEngine\n10x Faster"
    et.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in et.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Step 3: AI Processing (Right of engine)
    ai = slide.shapes.add_shape(
        MSO_SHAPE.FLOWCHART_DECISION,
        Inches(5.2), Inches(2.4),
        Inches(1.5), Inches(1.2)
    )
    ai.fill.solid()
    ai.fill.fore_color.rgb = RGBColor(147, 51, 234)
    ai.line.width = Pt(2)
    ai.line.color.rgb = RGBColor(107, 33, 168)
    
    ait = ai.text_frame
    ait.text = "🤖\nAI\nAssist"
    ait.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in ait.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Arrow from engine to AI
    arrow_ai = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(5.05), Inches(2.85),
        Inches(0.25), Inches(0.3)
    )
    arrow_ai.fill.solid()
    arrow_ai.fill.fore_color.rgb = RGBColor(230, 0, 0)
    arrow_ai.line.fill.background()
    
    # Step 4: Output Results (Bottom - 3 outputs)
    output_y = 4.5
    outputs = [
        ("📊\nExcel\nReports", RGBColor(255, 235, 156), 1.0),
        ("🌐\nHTML\nDiffs", RGBColor(179, 229, 252), 3.5),
        ("💬\nAI\nInsights", RGBColor(206, 237, 199), 6.0)
    ]
    
    for text, color, x in outputs:
        out = slide.shapes.add_shape(
            MSO_SHAPE.FLOWCHART_DOCUMENT,
            Inches(x), Inches(output_y),
            Inches(1.3), Inches(1.2)
        )
        out.fill.solid()
        out.fill.fore_color.rgb = color
        out.line.width = Pt(2)
        out.line.color.rgb = RGBColor(0, 102, 51)
        
        ot = out.text_frame
        ot.text = text
        ot.vertical_anchor = MSO_ANCHOR.TOP
        for p in ot.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(13)
            p.font.bold = True
    
    # Arrows from engine/AI to outputs
    down_arrows = [
        (3.5, 3.8, 1.5, 0.6),
        (4.2, 3.8, 3.9, 0.6),
        (5.5, 3.7, 6.3, 0.7)
    ]
    
    for x1, y1, x2, y2 in down_arrows:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            Inches(x2), Inches(y1),
            Inches(0.4), Inches(y2)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0, 153, 51)
        arrow.line.fill.background()
    
    # Bottom stats strip
    stats_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(6.2),
        Inches(9), Inches(1)
    )
    stats_bg.fill.solid()
    stats_bg.fill.fore_color.rgb = RGBColor(255, 245, 230)
    stats_bg.line.width = Pt(2)
    stats_bg.line.color.rgb = RGBColor(255, 140, 0)
    
    stats_text = stats_bg.text_frame
    stats_text.text = "✨ Features: Multi-Mode Comparison • RTC/ALM Integration • Parallel Processing • AI Merge • Real-time Progress"
    st = stats_text.paragraphs[0]
    st.font.size = Pt(15)
    st.font.bold = True
    st.alignment = PP_ALIGN.CENTER
    st.font.color.rgb = RGBColor(102, 51, 0)
    stats_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide

def add_time_savings_mega_slide(prs):
    """Slide 3: Dramatic time savings visualization"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Gradient background - dramatic
    bg1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(3.75))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = RGBColor(255, 250, 240)
    bg1.line.fill.background()
    
    bg2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.75), Inches(10), Inches(3.75))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = RGBColor(240, 255, 240)
    bg2.line.fill.background()
    
    # Red strip
    red_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.3))
    red_strip.fill.solid()
    red_strip.fill.fore_color.rgb = RGBColor(230, 0, 0)
    red_strip.line.fill.background()
    
    # Title with impact
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(9), Inches(0.8))
    tf = title.text_frame
    tf.text = "⚡ 10x FASTER PERFORMANCE!"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(230, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    
    # Big speedometer/gauge visual concept using shapes
    gauge_center_x = 5
    gauge_center_y = 3.2
    
    # Gauge background (semi-circle effect with rectangles)
    gauge_bg = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(2.5), Inches(1.8),
        Inches(5), Inches(3)
    )
    gauge_bg.fill.solid()
    gauge_bg.fill.fore_color.rgb = RGBColor(50, 50, 50)
    gauge_bg.line.width = Pt(4)
    gauge_bg.line.color.rgb = RGBColor(255, 140, 0)
    
    # Inner gauge
    gauge_inner = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(3), Inches(2.3),
        Inches(4), Inches(2.4)
    )
    gauge_inner.fill.solid()
    gauge_inner.fill.fore_color.rgb = RGBColor(20, 20, 20)
    gauge_inner.line.fill.background()
    
    # Speed indicator text
    speed_text = slide.shapes.add_textbox(Inches(3.5), Inches(2.7), Inches(3), Inches(1))
    stf = speed_text.text_frame
    stf.text = "10X\nSPEED"
    stf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in stf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 215, 0)
    
    # Comparison bars below gauge
    bar_y = 5.0
    
    # BEFORE section
    before_label = slide.shapes.add_textbox(Inches(0.5), Inches(bar_y), Inches(2), Inches(0.4))
    blt = before_label.text_frame
    blt.text = "❌ BEFORE"
    blt.paragraphs[0].font.size = Pt(18)
    blt.paragraphs[0].font.bold = True
    blt.paragraphs[0].font.color.rgb = RGBColor(200, 0, 0)
    
    # Long bar (before)
    before_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(bar_y + 0.5),
        Inches(8), Inches(0.5)
    )
    before_bar.fill.solid()
    before_bar.fill.fore_color.rgb = RGBColor(255, 100, 100)
    before_bar.line.width = Pt(2)
    before_bar.line.color.rgb = RGBColor(200, 0, 0)
    
    bbt = before_bar.text_frame
    bbt.text = "200 files: 10-15 minutes"
    bbt.paragraphs[0].font.size = Pt(18)
    bbt.paragraphs[0].font.bold = True
    bbt.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    bbt.paragraphs[0].alignment = PP_ALIGN.CENTER
    bbt.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # AFTER section
    after_label = slide.shapes.add_textbox(Inches(0.5), Inches(bar_y + 1.2), Inches(2), Inches(0.4))
    alt = after_label.text_frame
    alt.text = "✅ AFTER"
    alt.paragraphs[0].font.size = Pt(18)
    alt.paragraphs[0].font.bold = True
    alt.paragraphs[0].font.color.rgb = RGBColor(0, 153, 0)
    
    # Short bar (after)
    after_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(bar_y + 1.7),
        Inches(1.2), Inches(0.5)
    )
    after_bar.fill.solid()
    after_bar.fill.fore_color.rgb = RGBColor(100, 255, 100)
    after_bar.line.width = Pt(2)
    after_bar.line.color.rgb = RGBColor(0, 153, 0)
    
    abt = after_bar.text_frame
    abt.text = "1-2 minutes!"
    abt.paragraphs[0].font.size = Pt(18)
    abt.paragraphs[0].font.bold = True
    abt.paragraphs[0].font.color.rgb = RGBColor(0, 102, 0)
    abt.paragraphs[0].alignment = PP_ALIGN.CENTER
    abt.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Benefit badges on the right
    benefits = [
        ("🚀 Parallel\nProcessing", 3.2),
        ("💾 Memory\nOptimized", 4.5),
        ("📏 Smart\nFiltering", 5.8),
        ("⚡ Real-Time\nUpdates", 7.1)
    ]
    
    badge_colors = [
        RGBColor(200, 230, 255),
        RGBColor(255, 230, 200),
        RGBColor(220, 255, 220),
        RGBColor(255, 220, 255)
    ]
    
    for i, (text, x) in enumerate(benefits):
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(bar_y + 1.2),
            Inches(1.2), Inches(1)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = badge_colors[i]
        badge.line.width = Pt(2)
        badge.line.color.rgb = RGBColor(0, 51, 102)
        
        bt = badge.text_frame
        bt.text = text
        bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in bt.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(13)
            p.font.bold = True
    
    return slide

def add_features_circular_slide(prs):
    """Slide 4: Features in circular/hexagonal layout"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(245, 250, 255)
    bg.line.fill.background()
    
    # Red strip
    red_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.3))
    red_strip.fill.solid()
    red_strip.fill.fore_color.rgb = RGBColor(230, 0, 0)
    red_strip.line.fill.background()
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(9), Inches(0.7))
    tf = title.text_frame
    tf.text = "✨ Complete Feature Set"
    p = tf.paragraphs[0]
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER
    
    # Central hub
    hub = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON,
        Inches(3.8), Inches(3),
        Inches(2.4), Inches(2)
    )
    hub.fill.solid()
    hub.fill.fore_color.rgb = RGBColor(0, 51, 102)
    hub.line.width = Pt(4)
    hub.line.color.rgb = RGBColor(230, 0, 0)
    
    ht = hub.text_frame
    ht.text = "Migration\nAnalysis\nTool"
    ht.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in ht.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Features around the hub (circular layout)
    features = [
        ("🔍\nInterface\nAnalysis", 1.2, 1.8, RGBColor(200, 230, 255)),
        ("📊\nHTML/Excel\nReports", 4.2, 1.5, RGBColor(255, 230, 200)),
        ("☁️\nRTC/ALM\nIntegration", 7.2, 2.2, RGBColor(220, 255, 220)),
        ("⚡\n10x Faster\nProcessing", 7.8, 4.8, RGBColor(255, 220, 255)),
        ("🤖\nAI Smart\nMerge", 4.6, 5.8, RGBColor(255, 255, 200)),
        ("🔗\nDependency\nMapping", 1.0, 4.5, RGBColor(255, 200, 200))
    ]
    
    # Draw connecting lines first
    connections = [
        (2.3, 2.8, 3.8, 3.5),
        (5.3, 2.5, 4.7, 3.2),
        (8.1, 3.2, 6.2, 3.7),
        (8.2, 5.3, 6.2, 4.5),
        (5.7, 6.5, 5.3, 5.0),
        (2.1, 5.5, 3.8, 4.7)
    ]
    
    for x1, y1, x2, y2 in connections:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(x1), Inches(y1),
            Inches(abs(x2 - x1)), Inches(0.15)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(200, 200, 200)
        line.line.fill.background()
        line.rotation = 0
    
    # Draw feature boxes
    for text, x, y, color in features:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(1.5), Inches(1.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.width = Pt(3)
        box.line.color.rgb = RGBColor(0, 51, 102)
        box.shadow.inherit = False
        
        bt = box.text_frame
        bt.text = text
        bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in bt.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(12)
            p.font.bold = True
    
    return slide

def add_getting_started_flow_slide(prs):
    """Slide 5: Getting started with step-by-step flow"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background with two-tone
    bg_top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.3))
    bg_top.fill.solid()
    bg_top.fill.fore_color.rgb = RGBColor(0, 51, 102)
    bg_top.line.fill.background()
    
    bg_bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.3), Inches(10), Inches(6.2))
    bg_bottom.fill.solid()
    bg_bottom.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg_bottom.line.fill.background()
    
    # Red strip
    red_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    red_strip.fill.solid()
    red_strip.fill.fore_color.rgb = RGBColor(230, 0, 0)
    red_strip.line.fill.background()
    
    # Title in white on dark background
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.9))
    tf = title.text_frame
    tf.text = "🚀 Get Started in 4 Easy Steps"
    p = tf.paragraphs[0]
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Vertical flow with numbered steps
    steps = [
        ("1", "Launch Application", "python main.py or .exe", RGBColor(200, 230, 255), "💻"),
        ("2", "Select Mode", "Offline / Online / Hybrid", RGBColor(255, 230, 200), "⚙️"),
        ("3", "Configure & Run", "Set sources, RTC login", RGBColor(220, 255, 220), "🔧"),
        ("4", "View Results", "HTML/Excel reports + AI chat", RGBColor(255, 220, 255), "📊")
    ]
    
    step_y = 1.8
    
    for num, title, desc, color, icon in steps:
        # Number circle
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), Inches(step_y),
            Inches(0.8), Inches(0.8)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = RGBColor(230, 0, 0)
        num_circle.line.width = Pt(3)
        num_circle.line.color.rgb = RGBColor(0, 51, 102)
        
        nt = num_circle.text_frame
        nt.text = num
        nt.paragraphs[0].font.size = Pt(28)
        nt.paragraphs[0].font.bold = True
        nt.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        nt.paragraphs[0].alignment = PP_ALIGN.CENTER
        nt.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Step box
        step_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(step_y - 0.1),
            Inches(7.5), Inches(1)
        )
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = color
        step_box.line.width = Pt(2)
        step_box.line.color.rgb = RGBColor(0, 51, 102)
        
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(2.2), Inches(step_y + 0.1), Inches(0.7), Inches(0.7))
        it = icon_box.text_frame
        it.text = icon
        it.paragraphs[0].font.size = Pt(40)
        it.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Title and description
        text_box = slide.shapes.add_textbox(Inches(3.1), Inches(step_y), Inches(6), Inches(0.9))
        tt = text_box.text_frame
        tt.text = f"{title}\n{desc}"
        tt.paragraphs[0].font.size = Pt(20)
        tt.paragraphs[0].font.bold = True
        tt.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        tt.paragraphs[1].font.size = Pt(14)
        tt.paragraphs[1].font.color.rgb = RGBColor(51, 51, 51)
        tt.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Arrow between steps (except last)
        if num != "4":
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(1.1), Inches(step_y + 0.9),
                Inches(0.3), Inches(0.6)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(230, 0, 0)
            arrow.line.fill.background()
        
        step_y += 1.3
    
    # Bottom success banner
    banner = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(6.8),
        Inches(8), Inches(0.6)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(255, 215, 0)
    banner.line.width = Pt(3)
    banner.line.color.rgb = RGBColor(255, 140, 0)
    
    bt = banner.text_frame
    bt.text = "🎉 Start Migrating Faster Today! Contact: Bosch Engineering Team"
    bt.paragraphs[0].font.size = Pt(18)
    bt.paragraphs[0].font.bold = True
    bt.paragraphs[0].alignment = PP_ALIGN.CENTER
    bt.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide

def create_compact_presentation():
    """Create the compact 5-slide visual presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print("\n" + "=" * 75)
    print("  CREATING COMPACT 5-SLIDE VISUAL PRESENTATION")
    print("=" * 75)
    print()
    
    # Slide 1: Stunning Title
    add_stunning_title_slide(prs)
    print("  ✓ Slide 1: Stunning title with decorative elements")
    
    # Slide 2: Overview with complete workflow
    add_overview_flow_slide(prs)
    print("  ✓ Slide 2: Complete workflow diagram (Sources → Engine → AI → Results)")
    
    # Slide 3: Time Savings (BIG visual impact)
    add_time_savings_mega_slide(prs)
    print("  ✓ Slide 3: Dramatic time savings with gauge and comparison bars")
    
    # Slide 4: Features in circular layout
    add_features_circular_slide(prs)
    print("  ✓ Slide 4: Feature hub with circular connections")
    
    # Slide 5: Getting Started Flow
    add_getting_started_flow_slide(prs)
    print("  ✓ Slide 5: Step-by-step getting started flow")
    
    return prs

if __name__ == "__main__":
    print()
    prs = create_compact_presentation()
    
    output_path = "Migration_Tool_5Slides_Visual.pptx"
    prs.save(output_path)
    
    print()
    print("=" * 75)
    print("  ✨ PRESENTATION CREATED SUCCESSFULLY! ✨")
    print("=" * 75)
    print(f"\n📁 File: {output_path}")
    print(f"📊 Slides: {len(prs.slides)} (Compact & Powerful)")
    print("\n🎨 VISUAL ELEMENTS INCLUDED:")
    print("  ✓ Stunning title slide with decorative circles")
    print("  ✓ Complete workflow diagram with arrows")
    print("  ✓ Gauge visualization for speed")
    print("  ✓ Before/After comparison bars")
    print("  ✓ Circular feature hub with connections")
    print("  ✓ Step-by-step numbered flow")
    print("  ✓ Color-coded information boxes")
    print("  ✓ Professional Bosch branding")
    print("  ✓ 40+ visual shapes and diagrams")
    print("\n⚡ KEY HIGHLIGHTS:")
    print("  • Slide 2: Full workflow (Input → Processing → AI → Output)")
    print("  • Slide 3: 10x speed boost visualization")
    print("  • Slide 4: Feature ecosystem diagram")
    print("  • Slide 5: Getting started in 4 steps")
    print("\n🎯 Perfect for quick presentations (5-8 minutes)")
    print("=" * 75)
    print()
