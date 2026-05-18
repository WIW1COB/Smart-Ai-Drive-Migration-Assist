"""
Create an enhanced PowerPoint presentation with images, diagrams, and visual elements
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

def add_visual_title_slide(prs, title, subtitle):
    """Add an enhanced title slide with visual elements"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background gradient effect using shapes
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(234, 243, 251)
    bg_shape.line.fill.background()
    
    # Red strip at top (Bosch red)
    red_strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.3)
    )
    red_strip.fill.solid()
    red_strip.fill.fore_color.rgb = RGBColor(230, 0, 0)
    red_strip.line.color.rgb = RGBColor(230, 0, 0)
    
    # Decorative circles (abstract design)
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(8), Inches(5.5),
        Inches(2.5), Inches(2.5)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RGBColor(0, 51, 102)
    circle1.fill.transparency = 0.8
    circle1.line.fill.background()
    
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-0.5), Inches(0.5),
        Inches(2), Inches(2)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(230, 0, 0)
    circle2.fill.transparency = 0.7
    circle2.line.fill.background()
    
    # Main icon - Computer/Tool representation
    icon_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4), Inches(1.8),
        Inches(2), Inches(1.5)
    )
    icon_box.fill.solid()
    icon_box.fill.fore_color.rgb = RGBColor(0, 51, 102)
    icon_box.line.width = Pt(3)
    icon_box.line.color.rgb = RGBColor(230, 0, 0)
    
    icon_text = icon_box.text_frame
    icon_text.text = "⚙️🔧"
    p = icon_text.paragraphs[0]
    p.font.size = Pt(60)
    p.alignment = PP_ALIGN.CENTER
    icon_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.7), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(102, 102, 102)
    p.alignment = PP_ALIGN.CENTER
    
    # Bottom decorative bar
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.2),
        Inches(10), Inches(0.3)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = RGBColor(0, 51, 102)
    bottom_bar.line.fill.background()
    
    return slide

def add_visual_content_slide(prs, title, content_items, icon=""):
    """Add a visually enhanced content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg_shape.line.fill.background()
    
    # Red strip at top
    red_strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.25)
    )
    red_strip.fill.solid()
    red_strip.fill.fore_color.rgb = RGBColor(230, 0, 0)
    red_strip.line.color.rgb = RGBColor(230, 0, 0)
    
    # Side accent bar
    side_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.25),
        Inches(0.15), Inches(7.25)
    )
    side_bar.fill.solid()
    side_bar.fill.fore_color.rgb = RGBColor(0, 51, 102)
    side_bar.line.fill.background()
    
    # Title background
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.15), Inches(0.4),
        Inches(9.85), Inches(1)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(234, 243, 251)
    title_bg.line.fill.background()
    
    # Title with icon
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = f"{icon} {title}" if icon else title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Content with visual bullets
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(5.2))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.space_before = Pt(14)
        p.level = 0
        
        # Sub-bullets
        if item.strip().startswith('-'):
            p.text = item.replace('-', '', 1).strip()
            p.level = 1
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(51, 51, 51)
    
    return slide

def add_comparison_modes_slide(prs):
    """Visual slide showing comparison modes with diagrams"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg_shape.line.fill.background()
    
    # Red strip
    red_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.25))
    red_strip.fill.solid()
    red_strip.fill.fore_color.rgb = RGBColor(230, 0, 0)
    red_strip.line.color.rgb = RGBColor(230, 0, 0)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "📊 Comparison Modes"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Mode 1: Offline → Offline
    mode1_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(2.8), Inches(1.5))
    mode1_box.fill.solid()
    mode1_box.fill.fore_color.rgb = RGBColor(200, 230, 255)
    mode1_box.line.width = Pt(2)
    mode1_box.line.color.rgb = RGBColor(0, 51, 102)
    
    mode1_text = mode1_box.text_frame
    mode1_text.text = "📁 Offline → Offline\n\nFolder/ZIP\nComparison"
    for p in mode1_text.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)
    mode1_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Mode 2: Online → Online
    mode2_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.6), Inches(1.8), Inches(2.8), Inches(1.5))
    mode2_box.fill.solid()
    mode2_box.fill.fore_color.rgb = RGBColor(255, 230, 200)
    mode2_box.line.width = Pt(2)
    mode2_box.line.color.rgb = RGBColor(230, 0, 0)
    
    mode2_text = mode2_box.text_frame
    mode2_text.text = "☁️ Online → Online\n\nRTC Snapshot\nComparison"
    for p in mode2_text.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(102, 51, 0)
    mode2_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Mode 3: Hybrid
    mode3_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.8), Inches(2.8), Inches(1.5))
    mode3_box.fill.solid()
    mode3_box.fill.fore_color.rgb = RGBColor(220, 255, 220)
    mode3_box.line.width = Pt(2)
    mode3_box.line.color.rgb = RGBColor(0, 102, 51)
    
    mode3_text = mode3_box.text_frame
    mode3_text.text = "🔄 Online → Offline\n\nHybrid\nComparison"
    for p in mode3_text.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 51)
    mode3_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Features for each mode
    features = [
        ("Local directories\nZIP archives\nManual mapping", 0.5, 3.5),
        ("Live RTC snapshots\nChangeset tracking\nWork items", 3.6, 3.5),
        ("RTC + Local folder\nValidate changes\nFlexible workflow", 6.7, 3.5)
    ]
    
    for text, left, top in features:
        feature_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(2.8), Inches(1.2))
        tf = feature_box.text_frame
        tf.text = text
        for p in tf.paragraphs:
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(51, 51, 51)
    
    # Bottom visualization - Workflow diagram
    workflow_y = 5.2
    
    # Source
    source = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(workflow_y), Inches(1.5), Inches(0.8))
    source.fill.solid()
    source.fill.fore_color.rgb = RGBColor(200, 230, 255)
    source.line.width = Pt(2)
    source.line.color.rgb = RGBColor(0, 51, 102)
    st = source.text_frame
    st.text = "Source A"
    st.paragraphs[0].alignment = PP_ALIGN.CENTER
    st.paragraphs[0].font.bold = True
    st.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Arrow 1
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.2), Inches(workflow_y + 0.2), Inches(1), Inches(0.4))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = RGBColor(230, 0, 0)
    arrow1.line.fill.background()
    
    # Tool
    tool = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.3), Inches(workflow_y - 0.2), Inches(1.4), Inches(1.2))
    tool.fill.solid()
    tool.fill.fore_color.rgb = RGBColor(0, 51, 102)
    tool.line.width = Pt(2)
    tool.line.color.rgb = RGBColor(230, 0, 0)
    tt = tool.text_frame
    tt.text = "⚙️\nAnalysis\nTool"
    for p in tt.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.font.size = Pt(11)
    tt.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Arrow 2
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.9), Inches(workflow_y + 0.2), Inches(1), Inches(0.4))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = RGBColor(230, 0, 0)
    arrow2.line.fill.background()
    
    # Target
    target = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(workflow_y), Inches(1.5), Inches(0.8))
    target.fill.solid()
    target.fill.fore_color.rgb = RGBColor(255, 230, 200)
    target.line.width = Pt(2)
    target.line.color.rgb = RGBColor(230, 0, 0)
    tt = target.text_frame
    tt.text = "Source B"
    tt.paragraphs[0].alignment = PP_ALIGN.CENTER
    tt.paragraphs[0].font.bold = True
    tt.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Result
    result = slide.shapes.add_shape(MSO_SHAPE.FLOWCHART_DOCUMENT, Inches(3.5), Inches(6.3), Inches(3), Inches(0.9))
    result.fill.solid()
    result.fill.fore_color.rgb = RGBColor(200, 255, 200)
    result.line.width = Pt(2)
    result.line.color.rgb = RGBColor(0, 153, 0)
    rt = result.text_frame
    rt.text = "📊 Analysis Report\n(HTML + Excel)"
    for p in rt.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 102, 0)
    rt.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide

def add_time_savings_visual_slide(prs):
    """Enhanced time savings slide with visual bars"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg_shape.line.fill.background()
    
    # Red strip
    red_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.25))
    red_strip.fill.solid()
    red_strip.fill.fore_color.rgb = RGBColor(230, 0, 0)
    red_strip.line.color.rgb = RGBColor(230, 0, 0)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "⚡ Time Savings: 10x Performance Boost!"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Visual comparison bars
    scenarios = [
        ("50 Files", "2-3 min", "15-20 sec", 2.0, 90),
        ("200 Files", "10-15 min", "1-2 min", 2.7, 87),
        ("Single Diff", "2-3 sec", "0.5-1 sec", 3.4, 67)
    ]
    
    for scenario, before, after, y_pos, reduction in scenarios:
        # Scenario label
        label = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(1.5), Inches(0.3))
        lt = label.text_frame
        lt.text = scenario
        lt.paragraphs[0].font.size = Pt(18)
        lt.paragraphs[0].font.bold = True
        lt.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # BEFORE bar (red/long)
        before_width = 5.0
        before_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.3), Inches(y_pos), Inches(before_width), Inches(0.3))
        before_bar.fill.solid()
        before_bar.fill.fore_color.rgb = RGBColor(255, 100, 100)
        before_bar.line.width = Pt(1)
        before_bar.line.color.rgb = RGBColor(200, 0, 0)
        bt = before_bar.text_frame
        bt.text = f"Before: {before}"
        bt.paragraphs[0].font.size = Pt(14)
        bt.paragraphs[0].font.bold = True
        bt.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # AFTER bar (green/short)
        after_width = before_width * (100 - reduction) / 100
        after_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.3), Inches(y_pos + 0.4), Inches(after_width), Inches(0.3))
        after_bar.fill.solid()
        after_bar.fill.fore_color.rgb = RGBColor(100, 255, 100)
        after_bar.line.width = Pt(1)
        after_bar.line.color.rgb = RGBColor(0, 153, 0)
        at = after_bar.text_frame
        at.text = f"After: {after}"
        at.paragraphs[0].font.size = Pt(14)
        at.paragraphs[0].font.bold = True
        at.paragraphs[0].font.color.rgb = RGBColor(0, 102, 0)
        at.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Speedup badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.6), Inches(y_pos + 0.1), Inches(0.9), Inches(0.6))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(255, 215, 0)
        badge.line.width = Pt(2)
        badge.line.color.rgb = RGBColor(255, 140, 0)
        badge_text = badge.text_frame
        badge_text.text = "10x ⚡"
        badge_text.paragraphs[0].font.size = Pt(16)
        badge_text.paragraphs[0].font.bold = True
        badge_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Bottom stats boxes
    stat_boxes = [
        ("🚀 Parallel\nProcessing", "10 files at once", 0.8, 5.5, RGBColor(200, 230, 255)),
        ("💾 Memory\nDiffs", "No temp files", 2.8, 5.5, RGBColor(255, 230, 200)),
        ("📏 Smart\nFiltering", "Skip large files", 4.8, 5.5, RGBColor(220, 255, 220)),
        ("📊 Real-Time\nProgress", "Live updates", 6.8, 5.5, RGBColor(255, 220, 255))
    ]
    
    for title, subtitle, left, top, color in stat_boxes:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(1.8), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.width = Pt(2)
        box.line.color.rgb = RGBColor(0, 51, 102)
        
        box_text = box.text_frame
        box_text.text = f"{title}\n\n{subtitle}"
        for p in box_text.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(12)
            p.font.bold = True
        box_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide

def add_features_visual_slide(prs):
    """Features slide with icons and visual elements"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background gradient
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
    bg_shape.line.fill.background()
    
    # Red strip
    red_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.25))
    red_strip.fill.solid()
    red_strip.fill.fore_color.rgb = RGBColor(230, 0, 0)
    red_strip.line.color.rgb = RGBColor(230, 0, 0)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "✨ Key Features"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Feature grid (3x3)
    features = [
        ("🔍", "Interface\nAnalysis", "Analyze header\nfile interfaces", 0.5, 1.8, RGBColor(200, 230, 255)),
        ("🔗", "Platform\nDependency", "Visualize\ndependencies", 3.6, 1.8, RGBColor(255, 230, 200)),
        ("🤖", "AI Smart\nMerge", "Intelligent\nsuggestions", 6.7, 1.8, RGBColor(220, 255, 220)),
        ("🔐", "RTC/ALM\nIntegration", "Live server\nconnectivity", 0.5, 4.0, RGBColor(255, 220, 255)),
        ("📊", "Comprehensive\nReporting", "HTML + Excel\nreports", 3.6, 4.0, RGBColor(255, 255, 200)),
        ("💬", "AI\nAssistant", "Ask questions\nget answers", 6.7, 4.0, RGBColor(200, 255, 255))
    ]
    
    for emoji, title, desc, left, top, color in features:
        # Feature box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.8), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.width = Pt(3)
        box.line.color.rgb = RGBColor(0, 51, 102)
        
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.1), Inches(2.4), Inches(0.5))
        it = icon_box.text_frame
        it.text = emoji
        it.paragraphs[0].font.size = Pt(48)
        it.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.7), Inches(2.4), Inches(0.5))
        tt = title_box.text_frame
        tt.text = title
        tt.paragraphs[0].font.size = Pt(16)
        tt.paragraphs[0].font.bold = True
        tt.paragraphs[0].alignment = PP_ALIGN.CENTER
        tt.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 1.2), Inches(2.4), Inches(0.5))
        dt = desc_box.text_frame
        dt.text = desc
        for p in dt.paragraphs:
            p.font.size = Pt(13)
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(51, 51, 51)
    
    return slide

def create_visual_presentation():
    """Create the enhanced visual presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print("Creating enhanced visual presentation...")
    
    # Slide 1: Visual Title
    add_visual_title_slide(prs, 
                          "Migration Analysis Tool",
                          "Automated Migration Analysis with AI & RTC Integration")
    print("  ✓ Slide 1: Visual title slide")
    
    # Slide 2: Overview with icon
    add_visual_content_slide(prs, "Tool Overview", [
        "🎯 Comprehensive solution for code migration analysis",
        "⚡ 10x faster processing with parallel execution",
        "🤖 AI-powered intelligent merge suggestions",
        "☁️ Enterprise RTC/ALM integration",
        "📊 Detailed HTML and Excel reporting",
        "💬 Interactive AI comparison assistant"
    ], "🎯")
    print("  ✓ Slide 2: Overview")
    
    # Slide 3: Comparison Modes (Visual)
    add_comparison_modes_slide(prs)
    print("  ✓ Slide 3: Comparison modes with diagrams")
    
    # Slide 4: Time Savings (Visual Bars)
    add_time_savings_visual_slide(prs)
    print("  ✓ Slide 4: Time savings with visual bars")
    
    # Slide 5: Features Grid (Visual)
    add_features_visual_slide(prs)
    print("  ✓ Slide 5: Features grid with icons")
    
    # Slide 6: RTC Integration
    add_visual_content_slide(prs, "RTC/ALM Integration", [
        "🔐 Direct RTC Server Connectivity",
        "- Secure authentication with credential caching",
        "- Support for multiple RTC workspaces",
        "",
        "📸 Snapshot Management",
        "- Fetch and compare live snapshots",
        "- Component-level analysis",
        "- Automatic changeset tracking",
        "",
        "📝 Work Item Integration",
        "- Link changes to work items",
        "- Fetch work item details",
        "- Track development progress"
    ], "🔐")
    print("  ✓ Slide 6: RTC Integration")
    
    # Slide 7: AI Features
    add_visual_content_slide(prs, "AI-Powered Features", [
        "🤖 Smart Merge Suggestions",
        "- Google Gemini AI integration",
        "- Context-aware conflict resolution",
        "- Intelligent code recommendations",
        "",
        "💬 AI Comparison Assistant",
        "- Ask natural language questions",
        "- Get detailed change explanations",
        "- Dependency and impact analysis",
        "",
        "🎯 Comment Detection",
        "- Identify comment-only changes",
        "- Filter noise from real changes",
        "- Improve analysis accuracy"
    ], "🤖")
    print("  ✓ Slide 7: AI Features")
    
    # Slide 8: Reporting
    add_visual_content_slide(prs, "Comprehensive Reporting", [
        "📊 Excel Reports",
        "- Detailed statistics and metrics",
        "- Component-wise breakdowns",
        "- Exportable for stakeholders",
        "",
        "🌐 HTML Diff Reports",
        "- Side-by-side color-coded diffs",
        "- Syntax highlighting",
        "- Interactive navigation",
        "",
        "📈 Visual Dashboards",
        "- Progress tracking",
        "- Change summaries",
        "- Interface compatibility checks"
    ], "📊")
    print("  ✓ Slide 8: Reporting")
    
    # Slide 9: Use Cases
    add_visual_content_slide(prs, "Real-World Use Cases", [
        "🚀 Platform Migration Projects",
        "- Migrate from old to new platforms",
        "- Identify breaking API changes",
        "- Track migration completeness",
        "",
        "✅ Release Validation",
        "- Verify release candidate quality",
        "- Compare against baselines",
        "- Ensure code integrity",
        "",
        "🔄 Merge Management",
        "- Resolve complex merge conflicts",
        "- AI-powered suggestions",
        "- Streamline integration workflows"
    ], "💼")
    print("  ✓ Slide 9: Use Cases")
    
    # Slide 10: Benefits
    add_visual_content_slide(prs, "Benefits Summary", [
        "⏱️ Massive Time Savings: 10x faster processing",
        "🎯 High Accuracy: Comprehensive change detection",
        "🤝 Team Collaboration: Shareable reports",
        "🔧 Easy to Use: Intuitive GUI interface",
        "🔐 Enterprise Ready: RTC/ALM integration",
        "🤖 Intelligent: AI-powered assistance",
        "📊 Complete Visibility: Detailed analytics",
        "💪 Scalable: Handle large codebases"
    ], "✅")
    print("  ✓ Slide 10: Benefits")
    
    # Slide 11: Tool Coverage
    add_visual_content_slide(prs, "Complete Tool Coverage", [
        "✓ File Comparison & Diff Generation",
        "✓ RTC Snapshot Management & Fetching",
        "✓ Changeset & Work Item Tracking",
        "✓ Interface Analysis (Header Files)",
        "✓ Platform Dependency Visualization",
        "✓ AI-Powered Merge Suggestions",
        "✓ Excel & HTML Report Generation",
        "✓ Parallel Processing Engine (10x faster)",
        "✓ Comment-Only Change Detection",
        "✓ ZIP Archive Support",
        "✓ Credential Management & Caching",
        "✓ Real-Time Progress Updates"
    ], "📋")
    print("  ✓ Slide 11: Tool Coverage")
    
    # Slide 12: Getting Started
    add_visual_content_slide(prs, "Getting Started - 4 Easy Steps", [
        "1️⃣ Launch the Application",
        "- Execute: python main.py",
        "- Or use the standalone executable",
        "",
        "2️⃣ Choose Comparison Mode",
        "- Offline (Folders/ZIPs)",
        "- Online (RTC Snapshots)",
        "- Hybrid (Online + Offline)",
        "",
        "3️⃣ Configure & Connect",
        "- Enter RTC credentials (if needed)",
        "- Select source and target",
        "- Enable desired features",
        "",
        "4️⃣ Analyze & Review",
        "- Run comparison",
        "- View HTML/Excel reports",
        "- Ask AI assistant questions"
    ], "🚀")
    print("  ✓ Slide 12: Getting Started")
    
    # Slide 13: Thank You
    add_visual_title_slide(prs,
                          "Thank You!",
                          "Questions? Let's Discuss!")
    print("  ✓ Slide 13: Thank You")
    
    return prs

if __name__ == "__main__":
    print("=" * 70)
    print("  CREATING ENHANCED VISUAL PRESENTATION")
    print("=" * 70)
    print()
    
    prs = create_visual_presentation()
    
    output_path = "Migration_Analysis_Tool_Visual_Presentation.pptx"
    prs.save(output_path)
    
    print()
    print("=" * 70)
    print("  ✓ PRESENTATION CREATED SUCCESSFULLY!")
    print("=" * 70)
    print(f"\nFile: {output_path}")
    print(f"Total Slides: {len(prs.slides)}")
    print("\n🎨 VISUAL ENHANCEMENTS INCLUDED:")
    print("  • Colorful background designs")
    print("  • Icon-based feature grid")
    print("  • Visual comparison bars (time savings)")
    print("  • Workflow diagrams with arrows")
    print("  • Color-coded comparison modes")
    print("  • Decorative shapes and accents")
    print("  • Professional Bosch branding")
    print("  • Enhanced typography and layout")
    print("\n📊 Ready to present!")
    print("=" * 70)
